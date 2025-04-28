import os
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a new presentation with 16:9 aspect ratio
prs = Presentation()
# Change slide size to 16:9 (default is 4:3)
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title slide
title_content_layout = prs.slide_layouts[1]  # Title and content
section_header_layout = prs.slide_layouts[2]  # Section header
two_content_layout = prs.slide_layouts[3]  # Two content
comparison_layout = prs.slide_layouts[4]  # Comparison
title_only_layout = prs.slide_layouts[5]  # Title only
blank_layout = prs.slide_layouts[6]  # Blank slide

# Global counter for figure numbering
fig_counter = 1

# Helper function to set font sizes in a text frame
def set_text_font_size(text_frame, size_pt=18):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size_pt)

# Helper function to add a title slide
def add_title_slide(title, subtitle="", notes=""):
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Set font sizes
    set_text_font_size(title_shape.text_frame, 32)
    set_text_font_size(subtitle_shape.text_frame, 18)
    
    # Add speaker notes if provided
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide

# Helper function to add a content slide
def add_content_slide(title, content_list=None, notes=""):
    slide = prs.slides.add_slide(title_content_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    if content_list:
        text_frame = content_shape.text_frame
        for i, item in enumerate(content_list):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = item
            p.level = 0
    
    # Set font sizes
    set_text_font_size(title_shape.text_frame, 32)
    set_text_font_size(content_shape.text_frame, 18)
    
    # Add speaker notes if provided
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    
    return slide, content_shape

# Helper function to add an image to a slide with a caption
def add_image_to_slide(slide, img_path, left, top, width=None, height=None, caption=None):
    global fig_counter
    
    if os.path.exists(img_path):
        # Add the image
        if width and height:
            pic = slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            pic = slide.shapes.add_picture(img_path, left, top)
            
        # Generate figure caption if not provided
        if caption is None:
            # Extract filename without extension and path
            base_filename = os.path.basename(img_path)
            filename_without_ext = os.path.splitext(base_filename)[0]
            # Convert snake_case to Title Case for better readability
            image_name = ' '.join(word.capitalize() for word in filename_without_ext.split('_'))
            caption = f"Fig. {fig_counter}: {image_name}"
        else:
            caption = f"Fig. {fig_counter}: {caption}"
            
        # Add caption below the image
        caption_left = left
        caption_top = top + pic.height
        caption_width = pic.width
        caption_height = Inches(0.3)
        
        textbox = slide.shapes.add_textbox(caption_left, caption_top, caption_width, caption_height)
        textbox.text = caption
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_text_font_size(textbox.text_frame, 12)
        
        # Increment figure counter
        fig_counter += 1
        
        return True
    return False

# Helper function to add a table to a slide
def add_table_to_slide(slide, data, left, top, width, height):
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Fill the table with data
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = str(data[i][j])
            # Set font size for table text
            set_text_font_size(cell.text_frame, 18)
            
    return table

# 1. Title Slide
title_notes = """
Today I'll be presenting our machine learning project on predicting breast cancer survival using gene expression data from The Cancer Genome Atlas. This project integrates high-dimensional genomic data analysis with interpretable machine learning techniques to identify prognostic biomarkers.
"""
title_slide = add_title_slide(
    "Predicting Breast Cancer Survival from Gene Expression Data",
    "An Interpretable Machine Learning Approach with TCGA-BRCA Dataset",
    title_notes
)

# 2. Background & Motivation Slide
background_content = [
    "• Breast cancer is the most common cancer in women worldwide",
    "• Genomic data provides valuable insights but is high-dimensional (challenge)",
    "• Survival prediction critical for treatment planning and personalized medicine",
    "• Need for interpretable models that provide biological insights",
    "• Opportunity to identify novel gene biomarkers"
]

background_notes = """
Breast cancer remains a significant health concern globally. While genomic data offers tremendous potential for personalized medicine, its high dimensionality presents analytical challenges. Our goal was to develop models that not only predict survival accurately but are also interpretable to clinicians and researchers. This interpretability is crucial for identifying meaningful biomarkers and understanding the biological mechanisms behind breast cancer progression.
"""

background_slide, content_shape = add_content_slide("Background & Motivation", background_content, background_notes)

# Add background images
add_image_to_slide(background_slide, "plots/cancer_gene_survival.png", Inches(7.5), Inches(1.5), width=Inches(4.5), 
                  caption="Cancer Gene Survival Patterns")
add_image_to_slide(background_slide, "plots/cancer_gene_distributions.png", Inches(7.5), Inches(4.0), width=Inches(4.5), 
                  caption="Cancer Gene Expression Distributions")

# 3. Dataset Overview Slide
dataset_content = [
    "• TCGA-BRCA: The Cancer Genome Atlas Breast Cancer dataset",
    "• Data types:",
    "  - RNA-Seq gene expression data (HTSeq counts)",
    "  - Clinical annotations including survival outcomes",
    "• Dataset size:",
    "  - 1,097 patients with primary breast cancer",
    "  - ~20,000 genes measured",
    "  - 140 clinical features per patient",
    "• Key outcomes: Overall survival time and vital status (event indicator)"
]

dataset_notes = """
Our analysis utilized the TCGA-BRCA dataset, which contains comprehensive RNA sequencing data measuring the expression of approximately 20,000 genes across 1,097 breast cancer patients. The dataset also includes detailed clinical annotations, allowing us to link gene expression patterns with survival outcomes. This rich dataset enables us to investigate the molecular basis of breast cancer survival while accounting for censored observations - patients who were still alive at their last follow-up.
"""

dataset_slide, content_shape = add_content_slide("Dataset Overview", dataset_content, dataset_notes)

# Add dataset visualizations
add_image_to_slide(dataset_slide, "plots/pca_analysis.png", Inches(7.5), Inches(1.5), width=Inches(4.5), 
                  caption="PCA Analysis of Gene Expression Data")
add_image_to_slide(dataset_slide, "plots/gene_variance_distribution.png", Inches(7.5), Inches(4.0), width=Inches(4.5), 
                  caption="Gene Variance Distribution")
add_image_to_slide(dataset_slide, "plots/cumulative_variance.png", Inches(1.5), Inches(4.0), width=Inches(4.5), 
                  caption="Cumulative Explained Variance")

# 4. Data Preprocessing I - Basic Processing
preprocessing1_content = [
    "• Data acquisition and merging:",
    "  - Combined gene expression with clinical data using patient IDs",
    "  - Matched 1,097 patients with expression and survival data",
    "• Expression normalization:",
    "  - Log2-transformation to stabilize variance",
    "  - Z-score normalization to make genes comparable",
    "• Dimensionality reduction:",
    "  - Filtered low-variance genes (noise reduction)",
    "  - Reduced from ~20,000 to ~5,000 informative genes"
]

preprocessing1_notes = """
Our preprocessing pipeline began with acquiring and merging the gene expression and clinical data. RNA sequencing data is inherently noisy and skewed, so we applied log2-transformation to stabilize variance across expression levels, followed by z-score normalization to make different genes comparable. The visualization shows how normalization improved the data distribution. To manage the high dimensionality, we filtered out genes with minimal variation across patients, reducing our feature space from approximately 20,000 to 5,000 genes while retaining most of the informative signal.
"""

preprocessing1_slide, content_shape = add_content_slide("Data Preprocessing I - Basic Processing", preprocessing1_content, preprocessing1_notes)

# Add preprocessing visualizations
add_image_to_slide(preprocessing1_slide, "plots/normalization_effect.png", Inches(7.5), Inches(1.5), width=Inches(4.5), 
                  caption="Effect of Log2 and Z-Score Normalization")
add_image_to_slide(preprocessing1_slide, "plots/feature_selection.png", Inches(1.5), Inches(3.5), width=Inches(4.5), 
                  caption="Feature Selection Based on Variance")

# 5. Data Preprocessing II - Advanced Techniques
preprocessing2_content = [
    "• Train-test split:",
    "  - 70% training, 30% testing",
    "  - Stratified by event occurrence",
    "• Class imbalance handling:",
    "  - ~80% censored vs ~20% events",
    "  - Addressed through weighted models",
    "• Outlier detection and handling:",
    "  - Identified expression outliers using z-scores",
    "  - Applied winsorization (capping) to extreme values",
    "• Multicollinearity analysis:",
    "  - Correlation analysis among genes",
    "  - Principal Component Analysis (PCA)"
]

preprocessing2_notes = """
We implemented several advanced preprocessing techniques to improve model robustness. After splitting our data into training and test sets (stratified by survival events), we addressed the significant class imbalance between censored and deceased patients. About 80% of our cohort remained alive at last follow-up, which could bias models toward predicting the majority class. 

We also detected and handled outliers in gene expression using winsorization, which caps extreme values rather than removing them. The plot shows the effect of this approach on the data distribution. Finally, we analyzed multicollinearity among genes, as gene expression data is highly correlated, which can destabilize models. These preprocessing steps were crucial for building reliable predictive models with this complex genomic dataset.
"""

preprocessing2_slide, content_shape = add_content_slide("Data Preprocessing II - Advanced Techniques", preprocessing2_content, preprocessing2_notes)

# Add advanced preprocessing visualizations
add_image_to_slide(preprocessing2_slide, "plots/outlier_capping_comparison.png", Inches(7.5), Inches(1.5), width=Inches(4.5), 
                  caption="Outlier Capping Comparison")
add_image_to_slide(preprocessing2_slide, "plots/class_imbalance.png", Inches(7.5), Inches(4.0), width=Inches(4.5), 
                  caption="Class Imbalance in Survival Events")
add_image_to_slide(preprocessing2_slide, "plots/train_test_survival_distribution.png", Inches(1.5), Inches(2.5), width=Inches(4.5), 
                  caption="Train-Test Survival Distribution")
add_image_to_slide(preprocessing2_slide, "plots/outlier_analysis.png", Inches(1.5), Inches(5.0), width=Inches(4.5), 
                  caption="Gene Expression Outlier Analysis")

# 6. Exploratory Data Analysis
eda_content = [
    "• Survival characteristics:",
    "  - Kaplan-Meier curve of overall cohort",
    "  - Median follow-up: ~3.8 years",
    "  - Censoring rate: ~80%",
    "• Gene expression patterns:",
    "  - PCA visualization of expression data",
    "  - Correlation heatmap of top variant genes",
    "  - Investigation of known cancer genes"
]

eda_notes = """
Our exploratory analysis revealed important characteristics of the cohort. The Kaplan-Meier curve shows the overall survival probability over time, with a relatively high censoring rate of about 80%. The median follow-up was approximately 3.8 years.

We also explored patterns in gene expression data using principal component analysis, which revealed natural clustering of patients that may correspond to different molecular subtypes of breast cancer. The PCA plot demonstrates how gene expression patterns separate patients into groups that might have different survival outcomes. Additionally, we investigated correlations among highly variable genes and examined the behavior of known cancer-related genes to establish baseline expectations for our machine learning models.
"""

eda_slide, content_shape = add_content_slide("Exploratory Data Analysis", eda_content, eda_notes)

# Add EDA visualizations
add_image_to_slide(eda_slide, "plots/kaplan_meier_overall.png", Inches(1.5), Inches(3.0), width=Inches(4.0), 
                  caption="Overall Kaplan-Meier Survival Curve")
add_image_to_slide(eda_slide, "plots/biomarker_pca.png", Inches(6.5), Inches(3.0), width=Inches(5.0), 
                  caption="PCA of Potential Biomarker Genes")
add_image_to_slide(eda_slide, "plots/kaplan_meier_train_test.png", Inches(1.5), Inches(5.5), width=Inches(4.0), 
                  caption="Train-Test Kaplan-Meier Comparison")
add_image_to_slide(eda_slide, "plots/survival_time_analysis.png", Inches(6.5), Inches(5.5), width=Inches(5.0), 
                  caption="Survival Time Distribution Analysis")

# Create additional EDA slide for more visualizations
eda_additional_slide = prs.slides.add_slide(title_only_layout)
eda_additional_slide.shapes.title.text = "Exploratory Data Analysis - Gene Expression Patterns"
set_text_font_size(eda_additional_slide.shapes.title.text_frame, 32)

eda_additional_notes = """
Here we can see additional exploratory visualizations of gene expression patterns. The correlation heatmap shows relationships between highly variable genes, which helps us understand gene co-expression patterns. The cancer gene boxplots show differential expression of known cancer-related genes across patient subgroups. These patterns informed our feature selection and modeling approaches.
"""
eda_additional_slide.notes_slide.notes_text_frame.text = eda_additional_notes

# Add additional EDA visualizations
add_image_to_slide(eda_additional_slide, "plots/biomarker_correlation_heatmap.png", Inches(1.0), Inches(1.5), width=Inches(5.5), 
                  caption="Biomarker Correlation Heatmap")
add_image_to_slide(eda_additional_slide, "plots/cancer_gene_boxplots.png", Inches(7.0), Inches(1.5), width=Inches(5.5), 
                  caption="Cancer Gene Expression Boxplots")
add_image_to_slide(eda_additional_slide, "plots/biomarker_variance_distribution.png", Inches(1.0), Inches(4.5), width=Inches(5.5), 
                  caption="Biomarker Variance Distribution")
add_image_to_slide(eda_additional_slide, "plots/biomarker_clusters.png", Inches(7.0), Inches(4.5), width=Inches(5.5), 
                  caption="Gene Expression Cluster Analysis")

# 7. Modeling Approach I - Cox Regression Models
cox_content = [
    "• Cox Proportional Hazards with regularization:",
    "  - Lasso (L1): Sparse feature selection (50 genes)",
    "  - Ridge (L2): Handles correlated genes, uses all features",
    "  - Elastic Net: Balanced approach (239 genes)",
    "• Hyperparameter tuning:",
    "  - 5-fold cross-validation",
    "  - Optimized for concordance index (C-index)",
    "• Interpretability: Coefficients directly relate to hazard ratio"
]

cox_notes = """
Our primary modeling approach used Cox Proportional Hazards models, which are specifically designed for survival data and can handle censored observations. We applied three different regularization techniques: Lasso for sparse feature selection, Ridge for handling correlated features, and Elastic Net which balances both approaches.

The Lasso model identified 50 genes with non-zero coefficients, while Elastic Net selected 239 genes. Ridge regression used all features but with different weights. We used 5-fold cross-validation to tune the regularization parameters, optimizing for the concordance index, which measures the model's ability to rank patients by risk. These models are highly interpretable because each coefficient directly relates to a gene's effect on survival hazard ratio.
"""

cox_slide, content_shape = add_content_slide("Modeling Approach I - Cox Regression Models", cox_content, cox_notes)

# Add Cox model visualizations
add_image_to_slide(cox_slide, "plots/lasso_selected_genes.png", Inches(7.0), Inches(1.5), width=Inches(5.0), 
                  caption="Lasso Selected Gene Coefficients")
add_image_to_slide(cox_slide, "plots/cox_forest_plot.png", Inches(7.0), Inches(4.5), width=Inches(5.0), 
                  caption="Cox Model Forest Plot")
add_image_to_slide(cox_slide, "plots/elasticnet_parameter_heatmap.png", Inches(1.5), Inches(4.5), width=Inches(4.5), 
                  caption="ElasticNet Parameter Tuning Heatmap")

# Create additional Cox models slide
cox_additional_slide = prs.slides.add_slide(title_only_layout)
cox_additional_slide.shapes.title.text = "Additional Cox Model Results"
set_text_font_size(cox_additional_slide.shapes.title.text_frame, 32)

cox_additional_notes = """
Here we show additional results from our Cox regression modeling. The Ridge model identified different sets of genes with varying coefficient weights, while the ElasticNet model provided a balance between Lasso and Ridge approaches. The coefficients plot shows the magnitude and direction of gene effects on survival risk, with positive coefficients indicating increased risk and negative coefficients indicating protective effects.
"""
cox_additional_slide.notes_slide.notes_text_frame.text = cox_additional_notes

# Add additional Cox model visualizations
add_image_to_slide(cox_additional_slide, "plots/ridge_top_genes.png", Inches(1.5), Inches(1.5), width=Inches(5.0), 
                  caption="Ridge Regression Top Genes")
add_image_to_slide(cox_additional_slide, "plots/features/elasticnet_coefficients.png", Inches(7.0), Inches(1.5), width=Inches(5.0), 
                  caption="ElasticNet Coefficient Magnitudes")
add_image_to_slide(cox_additional_slide, "plots/features/top_genes_heatmap.png", Inches(1.5), Inches(4.5), width=Inches(5.0), 
                  caption="Top Genes Expression Heatmap")
add_image_to_slide(cox_additional_slide, "plots/elasticnet_selected_genes.png", Inches(7.0), Inches(4.5), width=Inches(5.0), 
                  caption="ElasticNet Selected Genes")

# 8. Modeling Approach II - Tree-Based Models
tree_content = [
    "• Decision Tree for Survival:",
    "  - Adapted for survival data",
    "  - Highly interpretable decision rules",
    "  - Limited depth to prevent overfitting",
    "• Random Survival Forest:",
    "  - Ensemble of 100 survival trees",
    "  - Captures nonlinear gene interactions",
    "  - Higher predictive power but less interpretable",
    "  - Provides feature importance rankings"
]

tree_notes = """
To capture potential nonlinear relationships in gene expression data, we also implemented tree-based models. A decision tree for survival data provides easily interpretable rules for stratifying patients into risk groups. The visualization shows how the model creates decision paths based on gene expression thresholds.

We then extended this to a Random Survival Forest, an ensemble of 100 survival trees that captures complex interactions between genes. This model achieved our highest predictive performance but at some cost to interpretability. However, it still provides valuable information through feature importance rankings, identifying genes that most strongly influence survival prediction. These tree-based approaches complement our linear Cox models, potentially uncovering relationships that linear models might miss.
"""

tree_slide, content_shape = add_content_slide("Modeling Approach II - Tree-Based Models", tree_content, tree_notes)

# Add tree-based model visualizations
add_image_to_slide(tree_slide, "plots/decision_tree_visualization.png", Inches(1.0), Inches(3.5), width=Inches(5.0), 
                  caption="Decision Tree Structure Visualization")
add_image_to_slide(tree_slide, "plots/random_forest_feature_importance.png", Inches(7.0), Inches(3.5), width=Inches(5.0), 
                  caption="Random Forest Feature Importance")

# Create additional tree models slide
tree_additional_slide = prs.slides.add_slide(title_only_layout)
tree_additional_slide.shapes.title.text = "Additional Tree-Based Model Results"
set_text_font_size(tree_additional_slide.shapes.title.text_frame, 32)

tree_additional_notes = """
Additional results from our tree-based models are shown here. The decision tree feature importance ranking highlights which genes were most influential in tree splits. The confusion matrix shows the classification performance of our decision tree model. We also analyzed the overlap between features selected by linear models (Cox) and nonlinear models (Random Forest), finding both common and unique biomarkers identified by different modeling approaches.
"""
tree_additional_slide.notes_slide.notes_text_frame.text = tree_additional_notes

# Add additional tree model visualizations
add_image_to_slide(tree_additional_slide, "plots/decision_tree_feature_importance.png", Inches(1.5), Inches(1.5), width=Inches(5.0), 
                  caption="Decision Tree Feature Importance")
add_image_to_slide(tree_additional_slide, "plots/decision_tree_confusion_matrix.png", Inches(7.0), Inches(1.5), width=Inches(5.0), 
                  caption="Decision Tree Confusion Matrix")
add_image_to_slide(tree_additional_slide, "plots/features/linear_nonlinear_venn.png", Inches(4.0), Inches(4.0), width=Inches(5.0), 
                  caption="Linear vs. Nonlinear Model Feature Comparison")

# 9. Model Performance Comparison
comparison_slide = prs.slides.add_slide(title_only_layout)
comparison_slide.shapes.title.text = "Model Performance Comparison"

# Set title font size
set_text_font_size(comparison_slide.shapes.title.text_frame, 32)

comparison_notes = """
Comparing our models, we found that the Random Survival Forest achieved the highest performance with a C-index of 0.759, indicating it correctly ranks the survival times of pairs of patients about 76% of the time. The Cox-based models performed well too, with Elastic Net achieving a C-index of 0.749 and Lasso at 0.740.

There's a clear trade-off between model performance and interpretability. While Random Forest performed best, the Lasso model offers a good balance - strong predictive power with a concise set of 50 genes and high interpretability. For clinical applications, this balance is crucial, as interpretability enables biologists and clinicians to understand and trust the model's predictions.
"""

# Add speaker notes
comparison_slide.notes_slide.notes_text_frame.text = comparison_notes

# Create table data
table_data = [
    ["Model", "C-index", "Features", "Model Type", "Interpretability"],
    ["Random Forest", "0.759", "100", "Survival", "Medium"],
    ["Elastic Net", "0.749", "239", "Cox", "High"],
    ["Lasso", "0.740", "50", "Cox", "High"],
    ["Ridge", "0.715", "50", "Cox", "Medium"],
    ["Decision Tree", "0.667", "50", "Classification", "Very High"]
]

# Add the table
table = add_table_to_slide(comparison_slide, table_data, Inches(1), Inches(1.5), Inches(11), Inches(2.5))

# Format the table - highlight the header and best performing model
for j in range(len(table_data[0])):
    # Format header
    header_cell = table.cell(0, j)
    header_cell.fill.solid()
    header_cell.fill.fore_color.rgb = RGBColor(220, 230, 242)  # Light blue
    
    # Format best model row (Random Forest)
    best_cell = table.cell(1, j)
    best_cell.fill.solid()
    best_cell.fill.fore_color.rgb = RGBColor(230, 230, 250)  # Light purple

# Add visualization of performance comparison
add_image_to_slide(comparison_slide, "plots/evaluation/performance_comparison.png", Inches(1.5), Inches(4.2), width=Inches(5.0), 
                  caption="Performance Comparison Across Models")
add_image_to_slide(comparison_slide, "plots/model_cindex_comparison.png", Inches(7.0), Inches(4.2), width=Inches(5.0), 
                  caption="Model C-Index Comparison")

# 10. Key Gene Biomarkers
biomarker_content = [
    "• Top genes identified across models:",
    "  - Common genes between Lasso and Random Forest",
    "  - Overlap between different modeling approaches",
    "• Gene signature visualization:",
    "  - Coefficient plot showing effect direction and magnitude",
    "  - Color-coded by increased/decreased risk",
    "• Selected gene examples:",
    "  - Genes with strongest association to survival",
    "  - Positive/negative hazard ratio interpretation"
]

biomarker_notes = """
One of our key objectives was to identify gene biomarkers associated with breast cancer survival. The Venn diagram shows the overlap between gene signatures identified by different models. Notably, several genes were consistently selected across multiple modeling approaches, suggesting their robust association with survival outcomes.

The coefficient plot visualizes our gene signature from the Lasso model, with red bars indicating genes whose increased expression is associated with worse survival (higher hazard ratio) and green bars showing protective genes (lower hazard ratio). For example, gene ENSG00000266550 shows a positive coefficient, suggesting it may promote tumor progression, while gene ENSG00000281856 has a negative coefficient, potentially indicating a tumor suppressive role.

This gene signature not only provides prognostic value but also points to potential therapeutic targets for further investigation.
"""

biomarker_slide, content_shape = add_content_slide("Key Gene Biomarkers", biomarker_content, biomarker_notes)

# Add gene biomarker visualizations
add_image_to_slide(biomarker_slide, "plots/gene_signature_venn3.png", Inches(1.5), Inches(3.5), width=Inches(4.5), 
                  caption="Venn Diagram of Gene Signatures")
add_image_to_slide(biomarker_slide, "plots/elasticnet_selected_genes.png", Inches(7.0), Inches(3.5), width=Inches(4.5), 
                  caption="ElasticNet Selected Gene Coefficients")
add_image_to_slide(biomarker_slide, "plots/biomarker_variance_ratio.png", Inches(1.5), Inches(5.5), width=Inches(4.5), 
                  caption="Biomarker Variance Ratio")

# 11. Survival Stratification
stratification_content = [
    "• Risk group stratification:",
    "  - Patients divided into high/medium/low risk based on model predictions",
    "  - Kaplan-Meier curves showing separation between groups",
    "  - Log-rank test p-value showing statistical significance",
    "• Clinical relevance:",
    "  - Clear risk stratification enables personalized treatment decisions",
    "  - Potential to identify patients needing more aggressive intervention"
]

stratification_notes = """
A crucial test of our models' clinical utility is their ability to stratify patients into meaningful risk groups. Using our best-performing model, we divided patients into risk groups based on their predicted risk scores.

The Kaplan-Meier curves demonstrate clear separation between these groups, with a highly significant log-rank test p-value. The high-risk group showed substantially worse survival outcomes compared to the low-risk group, validating our model's ability to distinguish between patients with different prognoses.

This stratification has direct clinical relevance, as it could help clinicians identify patients who might benefit from more aggressive treatment approaches, while sparing low-risk patients from unnecessary interventions. This represents a step toward personalized medicine in breast cancer care.
"""

stratification_slide, content_shape = add_content_slide("Survival Stratification", stratification_content, stratification_notes)

# Add stratification visualizations
add_image_to_slide(stratification_slide, "plots/evaluation/observed_outcome_kaplan_meier.png", Inches(7.0), Inches(2.0), width=Inches(5.0), 
                  caption="Risk Group Stratification Kaplan-Meier")
add_image_to_slide(stratification_slide, "plots/biomarker_cluster_survival.png", Inches(1.5), Inches(3.5), width=Inches(4.5), 
                  caption="Biomarker Cluster Survival Analysis")

# 12. Biological Pathway Analysis
pathway_content = [
    "• Biological relevance of identified genes:",
    "  - DNA repair and cell cycle regulation pathways",
    "  - Apoptosis and programmed cell death",
    "  - Cell proliferation and tumor growth",
    "  - Immune response and inflammation",
    "• Pathway enrichment analysis:",
    "  - Overrepresented biological processes",
    "  - Connection to known cancer mechanisms"
]

pathway_notes = """
Moving beyond prediction, we analyzed the biological relevance of our identified gene signatures. Pathway analysis revealed enrichment in several cancer-related processes, including DNA repair, cell cycle regulation, apoptosis, and immune response pathways.

Many of the top genes in our signature have established roles in cancer biology. For instance, several genes involved in cell cycle checkpoints were identified, consistent with the known importance of cell cycle dysregulation in cancer progression. Additionally, genes involved in DNA damage response were prominent, reflecting breast cancer's sensitivity to genomic instability.

This biological interpretation closes the loop between statistical modeling and cancer biology, providing not just predictive markers but also insights into the mechanisms driving breast cancer outcomes.
"""

pathway_slide, content_shape = add_content_slide("Biological Pathway Analysis", pathway_content, pathway_notes)

# Add pathway analysis visualizations
add_image_to_slide(pathway_slide, "plots/biomarker_correlation_heatmap.png", Inches(7.0), Inches(2.0), width=Inches(5.0), 
                  caption="Biomarker Gene Correlation Heatmap")
add_image_to_slide(pathway_slide, "plots/gene_survival_heatmap.png", Inches(1.5), Inches(3.5), width=Inches(5.0), 
                  caption="Gene Expression vs. Survival Heatmap")

# 13. Challenges and Limitations
challenges_content = [
    "• High dimensionality: ~20,000 genes vs. ~1,000 patients",
    "• Class imbalance: 80% censored vs. 20% events",
    "• Multicollinearity: Highly correlated gene expression",
    "• Censoring complexities: Right-censored survival data",
    "• External validation: Need to test on independent cohorts"
]

challenges_notes = """
Our study faced several challenges inherent to genomic data analysis. The high dimensionality of gene expression data (20,000 genes) relative to sample size (1,000 patients) created a risk of overfitting. Although our preprocessing and regularization techniques mitigated this risk, it remains a limitation.

Class imbalance between censored and deceased patients challenged model training, potentially biasing toward the majority class. Additionally, gene expression data exhibits high multicollinearity, making it difficult to isolate individual gene effects.

Survival analysis itself introduces complexity through right-censoring, where the true survival time for many patients remains unknown. Finally, while we used cross-validation, external validation on independent cohorts is necessary to confirm our findings' generalizability before clinical application.
"""

challenges_slide, content_shape = add_content_slide("Challenges and Limitations", challenges_content, challenges_notes)

# Add challenges visualizations
add_image_to_slide(challenges_slide, "plots/evaluation/complexity_vs_performance.png", Inches(7.0), Inches(2.0), width=Inches(5.0), 
                  caption="Model Complexity vs. Performance Trade-off")
add_image_to_slide(challenges_slide, "plots/model_complexity_vs_performance.png", Inches(1.5), Inches(3.5), width=Inches(5.0), 
                  caption="Performance vs. Model Complexity")

# 14. Conclusions and Future Work
conclusion_content = [
    "• Conclusions:",
    "  - Successfully developed predictive models for breast cancer survival",
    "  - Random Forest achieved highest performance (C-index: 0.759)",
    "  - Identified robust gene signature with prognostic value",
    "  - Balanced model performance with biological interpretability",
    "• Future directions:",
    "  - External validation on independent datasets (e.g., METABRIC)",
    "  - Integration with clinical variables to improve prediction",
    "  - Experimental validation of key biomarkers",
    "  - Development of clinically applicable risk score"
]

conclusion_notes = """
In conclusion, we successfully developed machine learning models to predict breast cancer survival from gene expression data. Our Random Forest model achieved the highest predictive performance with a C-index of 0.759, while our Lasso Cox model offered a good balance of performance and interpretability.

We identified a robust gene signature with prognostic value, potentially offering new insights into breast cancer biology and progression. Our models effectively stratified patients into risk groups with significantly different survival outcomes, demonstrating clinical relevance.

Future work includes external validation on independent datasets, integration with clinical variables to improve predictive power, experimental validation of our identified biomarkers, and development of a clinically applicable risk score. This work represents a step toward more personalized treatment approaches for breast cancer patients based on their genomic profiles.
"""

conclusion_slide, content_shape = add_content_slide("Conclusions and Future Work", conclusion_content, conclusion_notes)

# Add conclusion visualizations
add_image_to_slide(conclusion_slide, "plots/evaluation/model_radar_comparison.png", Inches(7.0), Inches(2.0), width=Inches(5.0), 
                  caption="Model Performance Radar Comparison")

# 15. Thank You / Questions
thanks_notes = """
Thank you for your attention. I'd be happy to answer any questions about our methodology, findings, or potential implications for breast cancer research and treatment.
"""

thanks_slide = add_title_slide("Thank You!", "Questions?", thanks_notes)

# Save the presentation
prs.save("TCGA_BRCA_Survival_Prediction_Detailed_16_9.pptx")

print("Detailed presentation created successfully: TCGA_BRCA_Survival_Prediction_Detailed_16_9.pptx") 